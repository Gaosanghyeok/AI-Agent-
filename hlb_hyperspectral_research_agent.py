#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
HLB Hyperspectral Research Agent Pipeline
=========================================

用途：
    面向“基于高光谱图像的柑橘黄龙病叶片检测”的科研流程，构建一个可运行的多 Agent
    辅助系统。它不直接替代你的训练代码，而是把科研工程中的数据审计、预处理结果记录、
    5-fold OOF / 6:2:2 划分、指标汇总、混淆矩阵分析、组会汇报提纲和课程表单说明统一串起来。

适用场景：
    1. 证明你使用 Agent / AI 驱动构建了科研流水线。
    2. 帮你管理服务器上 500+ 片叶片数据的预处理、训练、验证和测试结果。
    3. 自动生成 Markdown/HTML 报告和 PPT 汇报提纲。

建议目录结构：
    project_root/
        manifest.csv                         # 样本清单，至少包含 label；推荐包含 sample_id, source_image, path
        raw_data/                            # 原始高光谱数据，可选
        processed_npz/                       # 预处理后 NPZ 或 ROI 数据，可选
        results/                             # 训练结果目录
            metrics_summary.csv              # 可选：模型指标汇总
            oof_confusion_matrix.csv         # 可选：5-fold OOF 混淆矩阵
            test_confusion_matrix.csv        # 可选：6:2:2 test 混淆矩阵
            test_confusion_matrix.png        # 可选：图片会复制到报告目录，但无法自动读取数值

运行示例：
    python hlb_hyperspectral_research_agent.py \
        --project_name "Citrus-HLB Hyperspectral Research Agent" \
        --manifest E:/HLB_project/manifest.csv \
        --raw_dir E:/HLB_project/raw_data \
        --processed_dir E:/HLB_project/processed_npz \
        --results_dir E:/HLB_project/results \
        --out_dir E:/HLB_project/agent_output \
        --label_col label \
        --group_col source_image \
        --make_splits

依赖：
    必需：Python 3.8+，numpy，pandas
    推荐：scikit-learn，matplotlib
    可选：python-pptx（如果希望自动生成简易 PPTX）

注意：
    - 如果你的混淆矩阵只有 png 图片，代码会复制图片到报告，但不能稳定读取数字；
      最好同时保存 csv/npy 版本。
    - 如果你的 manifest 中存在 source_image，建议使用 group_col=source_image，避免同一原始图像的叶片
      同时进入训练集和测试集，造成数据泄漏。
"""

from __future__ import annotations

import argparse
import csv
import dataclasses
import datetime as _dt
import json
import math
import os
import re
import shutil
import sys
from collections import Counter, defaultdict
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Sequence, Tuple

import numpy as np
import pandas as pd

try:
    import matplotlib.pyplot as plt
except Exception:  # pragma: no cover
    plt = None

try:
    from sklearn.model_selection import GroupKFold, GroupShuffleSplit, StratifiedKFold, train_test_split
    try:
        from sklearn.model_selection import StratifiedGroupKFold
    except Exception:  # pragma: no cover
        StratifiedGroupKFold = None
except Exception:  # pragma: no cover
    GroupKFold = None
    GroupShuffleSplit = None
    StratifiedKFold = None
    StratifiedGroupKFold = None
    train_test_split = None


DEFAULT_LABEL_ORDER = ["HA", "HE", "ND", "HS"]
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp"}
MATRIX_SUFFIXES = {".csv", ".npy", ".npz", ".txt"}


# -----------------------------
# 通用工具函数
# -----------------------------

def now_str() -> str:
    return _dt.datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def ensure_dir(path: Path) -> Path:
    path.mkdir(parents=True, exist_ok=True)
    return path


def safe_rel(path: Path, base: Path) -> str:
    try:
        return str(path.resolve().relative_to(base.resolve()))
    except Exception:
        return str(path)


def human_size(num_bytes: float) -> str:
    if num_bytes is None or (isinstance(num_bytes, float) and math.isnan(num_bytes)):
        return "NA"
    units = ["B", "KB", "MB", "GB", "TB", "PB"]
    size = float(num_bytes)
    for unit in units:
        if abs(size) < 1024.0:
            return f"{size:.2f} {unit}"
        size /= 1024.0
    return f"{size:.2f} EB"


def folder_size(path: Optional[Path]) -> int:
    if path is None or not path.exists():
        return 0
    total = 0
    for root, _, files in os.walk(path):
        for f in files:
            fp = Path(root) / f
            try:
                total += fp.stat().st_size
            except OSError:
                pass
    return total


def scan_files(path: Optional[Path], max_examples: int = 10) -> Dict[str, Any]:
    if path is None or not path.exists():
        return {"exists": False, "path": str(path) if path else None}
    ext_counter = Counter()
    n_files = 0
    examples = []
    total_size = 0
    for root, _, files in os.walk(path):
        for f in files:
            fp = Path(root) / f
            n_files += 1
            ext_counter[fp.suffix.lower() or "[no_suffix]"] += 1
            try:
                total_size += fp.stat().st_size
            except OSError:
                pass
            if len(examples) < max_examples:
                examples.append(str(fp))
    return {
        "exists": True,
        "path": str(path),
        "n_files": n_files,
        "total_size_bytes": total_size,
        "total_size_human": human_size(total_size),
        "extensions": dict(ext_counter.most_common()),
        "examples": examples,
    }


def read_csv_robust(path: Path) -> pd.DataFrame:
    encodings = ["utf-8-sig", "utf-8", "gbk", "gb18030"]
    last_err = None
    for enc in encodings:
        try:
            return pd.read_csv(path, encoding=enc)
        except Exception as e:
            last_err = e
    raise RuntimeError(f"无法读取 CSV: {path}; last error={last_err}")


def write_json(path: Path, obj: Any) -> None:
    with open(path, "w", encoding="utf-8") as f:
        json.dump(obj, f, ensure_ascii=False, indent=2)


def markdown_table(df: pd.DataFrame, max_rows: int = 30) -> str:
    if df is None or df.empty:
        return ""
    show = df.copy()
    if len(show) > max_rows:
        show = show.head(max_rows)
    # 转为字符串，避免 markdown 中出现 nan
    show = show.fillna("NA")
    cols = list(show.columns)
    lines = []
    lines.append("| " + " | ".join(map(str, cols)) + " |")
    lines.append("| " + " | ".join(["---"] * len(cols)) + " |")
    for _, row in show.iterrows():
        vals = [str(row[c]).replace("\n", " ") for c in cols]
        lines.append("| " + " | ".join(vals) + " |")
    if len(df) > max_rows:
        lines.append(f"\n> 仅显示前 {max_rows} 行，共 {len(df)} 行。")
    return "\n".join(lines)


# -----------------------------
# 数据结构
# -----------------------------

@dataclasses.dataclass
class AgentContext:
    project_name: str
    manifest_path: Optional[Path]
    raw_dir: Optional[Path]
    processed_dir: Optional[Path]
    results_dir: Optional[Path]
    out_dir: Path
    label_col: str = "label"
    group_col: Optional[str] = "source_image"
    id_col: Optional[str] = "sample_id"
    label_order: List[str] = dataclasses.field(default_factory=lambda: DEFAULT_LABEL_ORDER.copy())
    make_splits: bool = False

    manifest_df: Optional[pd.DataFrame] = None
    artifacts: Dict[str, Any] = dataclasses.field(default_factory=dict)
    logs: List[str] = dataclasses.field(default_factory=list)

    def log(self, msg: str) -> None:
        line = f"[{now_str()}] {msg}"
        print(line)
        self.logs.append(line)


# -----------------------------
# Base Agent
# -----------------------------

class BaseAgent:
    name = "BaseAgent"

    def run(self, ctx: AgentContext) -> None:
        raise NotImplementedError


# -----------------------------
# 1. 数据审计 Agent
# -----------------------------

class DataAuditAgent(BaseAgent):
    name = "DataAuditAgent"

    def run(self, ctx: AgentContext) -> None:
        ctx.log(f"{self.name}: 开始扫描数据与样本清单。")
        audit: Dict[str, Any] = {}
        audit["project_name"] = ctx.project_name
        audit["created_at"] = now_str()
        audit["raw_dir"] = scan_files(ctx.raw_dir)
        audit["processed_dir"] = scan_files(ctx.processed_dir)
        audit["results_dir"] = scan_files(ctx.results_dir)

        raw_size = audit["raw_dir"].get("total_size_bytes", 0) or 0
        processed_size = audit["processed_dir"].get("total_size_bytes", 0) or 0
        if raw_size > 0 and processed_size > 0:
            audit["compression_ratio_raw_to_processed"] = raw_size / processed_size
            audit["space_reduction_percent"] = (1 - processed_size / raw_size) * 100
        else:
            audit["compression_ratio_raw_to_processed"] = None
            audit["space_reduction_percent"] = None

        if ctx.manifest_path and ctx.manifest_path.exists():
            df = read_csv_robust(ctx.manifest_path)
            ctx.manifest_df = df
            audit["manifest"] = {
                "path": str(ctx.manifest_path),
                "n_rows": int(len(df)),
                "n_cols": int(df.shape[1]),
                "columns": list(df.columns),
            }
            if ctx.label_col in df.columns:
                vc = df[ctx.label_col].astype(str).value_counts().to_dict()
                audit["label_distribution"] = vc
            if ctx.group_col and ctx.group_col in df.columns:
                audit["n_groups"] = int(df[ctx.group_col].astype(str).nunique())
            if ctx.id_col and ctx.id_col in df.columns:
                audit["n_unique_ids"] = int(df[ctx.id_col].astype(str).nunique())
        else:
            audit["manifest"] = {"exists": False, "path": str(ctx.manifest_path) if ctx.manifest_path else None}
            ctx.log(f"{self.name}: 未提供或未找到 manifest.csv，后续划分功能会跳过。")

        ctx.artifacts["data_audit"] = audit
        write_json(ctx.out_dir / "data_audit.json", audit)
        ctx.log(f"{self.name}: 已生成 data_audit.json。")


# -----------------------------
# 2. 数据划分 Agent
# -----------------------------

class SplitDesignAgent(BaseAgent):
    name = "SplitDesignAgent"

    def run(self, ctx: AgentContext) -> None:
        if not ctx.make_splits:
            ctx.log(f"{self.name}: 未启用 --make_splits，跳过划分文件生成。")
            return
        df = ctx.manifest_df
        if df is None or df.empty:
            ctx.log(f"{self.name}: manifest 不存在，无法生成划分。")
            return
        if ctx.label_col not in df.columns:
            ctx.log(f"{self.name}: manifest 缺少 label_col={ctx.label_col}，无法生成划分。")
            return

        splits_dir = ensure_dir(ctx.out_dir / "splits")
        y = df[ctx.label_col].astype(str).values
        groups = None
        use_group = bool(ctx.group_col and ctx.group_col in df.columns)
        if use_group:
            groups = df[ctx.group_col].astype(str).values
            ctx.log(f"{self.name}: 使用 group_col={ctx.group_col} 做分组划分，降低数据泄漏风险。")
        else:
            ctx.log(f"{self.name}: 未找到 group_col，将按样本级 stratify 划分。")

        # 5-fold OOF
        fold_df = df.copy()
        fold_df["fold"] = -1
        try:
            if use_group and StratifiedGroupKFold is not None:
                splitter = StratifiedGroupKFold(n_splits=5, shuffle=True, random_state=42)
                iterator = splitter.split(df, y, groups)
            elif use_group and GroupKFold is not None:
                splitter = GroupKFold(n_splits=5)
                iterator = splitter.split(df, y, groups)
            elif StratifiedKFold is not None:
                splitter = StratifiedKFold(n_splits=5, shuffle=True, random_state=42)
                iterator = splitter.split(df, y)
            else:
                iterator = self._simple_kfold_indices(len(df), 5)

            for fold, (_, valid_idx) in enumerate(iterator):
                fold_df.loc[fold_df.index[valid_idx], "fold"] = fold
            fold_path = splits_dir / "split_5fold_oof.csv"
            fold_df.to_csv(fold_path, index=False, encoding="utf-8-sig")
            ctx.artifacts["split_5fold_oof"] = str(fold_path)
            ctx.log(f"{self.name}: 已生成 5-fold OOF 划分：{fold_path}")
        except Exception as e:
            ctx.log(f"{self.name}: 5-fold 划分失败：{repr(e)}")

        # 6:2:2 train/val/test
        try:
            split622 = self._make_622_split(df, ctx.label_col, ctx.group_col if use_group else None)
            split622_path = splits_dir / "split_6_2_2.csv"
            split622.to_csv(split622_path, index=False, encoding="utf-8-sig")
            ctx.artifacts["split_6_2_2"] = str(split622_path)
            ctx.log(f"{self.name}: 已生成 6:2:2 划分：{split622_path}")
        except Exception as e:
            ctx.log(f"{self.name}: 6:2:2 划分失败：{repr(e)}")

    @staticmethod
    def _simple_kfold_indices(n: int, k: int) -> Iterable[Tuple[np.ndarray, np.ndarray]]:
        idx = np.arange(n)
        np.random.default_rng(42).shuffle(idx)
        folds = np.array_split(idx, k)
        for i in range(k):
            valid = folds[i]
            train = np.setdiff1d(idx, valid)
            yield train, valid

    @staticmethod
    def _make_622_split(df: pd.DataFrame, label_col: str, group_col: Optional[str]) -> pd.DataFrame:
        out = df.copy()
        out["split"] = "train"

        if group_col and group_col in df.columns and GroupShuffleSplit is not None:
            # 先分 test=20%，再从剩余中分 val=25% => 总体约 20%
            groups = df[group_col].astype(str).values
            y = df[label_col].astype(str).values
            gss1 = GroupShuffleSplit(n_splits=1, test_size=0.2, random_state=42)
            train_val_idx, test_idx = next(gss1.split(df, y, groups))
            out.loc[out.index[test_idx], "split"] = "test"

            df_train_val = df.iloc[train_val_idx].copy()
            groups_tv = df_train_val[group_col].astype(str).values
            y_tv = df_train_val[label_col].astype(str).values
            gss2 = GroupShuffleSplit(n_splits=1, test_size=0.25, random_state=43)
            train_idx_local, val_idx_local = next(gss2.split(df_train_val, y_tv, groups_tv))
            val_idx = df_train_val.index[val_idx_local]
            out.loc[val_idx, "split"] = "val"
            return out

        if train_test_split is not None:
            idx = np.arange(len(df))
            y = df[label_col].astype(str).values
            train_val_idx, test_idx = train_test_split(
                idx, test_size=0.2, random_state=42, stratify=y
            )
            y_train_val = y[train_val_idx]
            train_idx, val_idx = train_test_split(
                train_val_idx, test_size=0.25, random_state=43, stratify=y_train_val
            )
            out.loc[out.index[test_idx], "split"] = "test"
            out.loc[out.index[val_idx], "split"] = "val"
            out.loc[out.index[train_idx], "split"] = "train"
            return out

        # 无 sklearn 的简单 fallback
        rng = np.random.default_rng(42)
        for label, sub in df.groupby(label_col):
            indices = sub.index.to_numpy()
            rng.shuffle(indices)
            n = len(indices)
            n_test = max(1, int(round(n * 0.2)))
            n_val = max(1, int(round(n * 0.2)))
            out.loc[indices[:n_test], "split"] = "test"
            out.loc[indices[n_test:n_test + n_val], "split"] = "val"
            out.loc[indices[n_test + n_val:], "split"] = "train"
        return out


# -----------------------------
# 3. 结果分析 Agent
# -----------------------------

class ResultAnalysisAgent(BaseAgent):
    name = "ResultAnalysisAgent"

    def run(self, ctx: AgentContext) -> None:
        ctx.log(f"{self.name}: 开始解析训练结果、指标表和混淆矩阵。")
        results: Dict[str, Any] = {
            "metrics_files": [],
            "confusion_matrices": [],
            "image_confusion_matrices": [],
            "analysis_summary": [],
        }
        if not ctx.results_dir or not ctx.results_dir.exists():
            ctx.log(f"{self.name}: results_dir 不存在，跳过结果分析。")
            ctx.artifacts["result_analysis"] = results
            return

        copied_images_dir = ensure_dir(ctx.out_dir / "figures")

        # 指标文件
        metric_candidates = self._find_metric_files(ctx.results_dir)
        for p in metric_candidates:
            parsed = self._parse_metric_file(p)
            if parsed:
                results["metrics_files"].append(parsed)

        # 混淆矩阵：csv/npy/npz/txt
        cm_candidates = self._find_confusion_files(ctx.results_dir)
        for p in cm_candidates:
            if p.suffix.lower() in IMAGE_SUFFIXES:
                dst = copied_images_dir / p.name
                try:
                    shutil.copy2(p, dst)
                except Exception:
                    pass
                results["image_confusion_matrices"].append({
                    "path": str(p),
                    "copied_to": str(dst),
                    "note": "图片版混淆矩阵已复制，但建议同时保存 CSV/NPY 版本以便自动计算指标。",
                })
                continue
            cm = self._load_confusion_matrix(p, ctx.label_order)
            if cm is not None:
                cm_report = self._analyze_cm(cm, ctx.label_order, matrix_name=p.name)
                cm_report["source_path"] = str(p)
                results["confusion_matrices"].append(cm_report)

        # 自动总结
        summaries = []
        if results["metrics_files"]:
            summaries.append(self._summarize_metrics(results["metrics_files"]))
        for cm_report in results["confusion_matrices"]:
            summaries.append(cm_report.get("text_summary", ""))
        if results["image_confusion_matrices"] and not results["confusion_matrices"]:
            summaries.append("检测到图片版混淆矩阵，但未检测到可解析的 CSV/NPY 矩阵。建议训练脚本保存 test_confusion_matrix.csv 和 oof_confusion_matrix.csv。")
        results["analysis_summary"] = [s for s in summaries if s]

        write_json(ctx.out_dir / "result_analysis.json", results)
        ctx.artifacts["result_analysis"] = results
        ctx.log(f"{self.name}: 已生成 result_analysis.json。")

    @staticmethod
    def _find_metric_files(root: Path) -> List[Path]:
        keys = ["metric", "summary", "result", "classification_report", "score"]
        files = []
        for p in root.rglob("*"):
            if not p.is_file():
                continue
            name = p.name.lower()
            if any(k in name for k in keys) and p.suffix.lower() in {".csv", ".json", ".txt", ".log"}:
                files.append(p)
        return sorted(files)

    @staticmethod
    def _find_confusion_files(root: Path) -> List[Path]:
        files = []
        for p in root.rglob("*"):
            if not p.is_file():
                continue
            name = p.name.lower()
            if "confusion" in name or "cm" in name:
                if p.suffix.lower() in MATRIX_SUFFIXES or p.suffix.lower() in IMAGE_SUFFIXES:
                    files.append(p)
        return sorted(files)

    @staticmethod
    def _parse_metric_file(path: Path) -> Optional[Dict[str, Any]]:
        suffix = path.suffix.lower()
        try:
            if suffix == ".csv":
                df = read_csv_robust(path)
                # 只保留常见指标列，若没有也保留前几列
                preferred = [c for c in df.columns if re.search(r"acc|f1|precision|recall|kappa|mcc|auc|loss|model|split|fold", str(c), re.I)]
                preview_cols = preferred if preferred else list(df.columns[:12])
                preview = df[preview_cols].head(20).to_dict(orient="records")
                best_row = None
                score_col = None
                for c in df.columns:
                    lc = str(c).lower()
                    if lc in {"macro_f1", "macro-f1", "f1_macro"} or "macro" in lc and "f1" in lc:
                        score_col = c
                        break
                if score_col is None:
                    for c in df.columns:
                        if "accuracy" in str(c).lower() or str(c).lower() == "acc":
                            score_col = c
                            break
                if score_col is not None:
                    tmp = pd.to_numeric(df[score_col], errors="coerce")
                    if tmp.notna().any():
                        best_idx = int(tmp.idxmax())
                        best_row = df.loc[best_idx].to_dict()
                return {
                    "path": str(path),
                    "type": "csv",
                    "n_rows": int(len(df)),
                    "columns": list(df.columns),
                    "preview": preview,
                    "best_by": score_col,
                    "best_row": best_row,
                }
            if suffix == ".json":
                with open(path, "r", encoding="utf-8") as f:
                    obj = json.load(f)
                flat = ResultAnalysisAgent._flatten_dict(obj)
                metric_keys = {k: v for k, v in flat.items() if re.search(r"acc|f1|precision|recall|kappa|mcc|loss", k, re.I)}
                return {"path": str(path), "type": "json", "metrics": metric_keys, "raw_keys": list(obj) if isinstance(obj, dict) else []}
            if suffix in {".txt", ".log"}:
                txt = path.read_text(encoding="utf-8", errors="ignore")[:20000]
                metrics = {}
                patterns = {
                    "accuracy": r"(?:accuracy|acc)\s*[:=]\s*([0-9]*\.?[0-9]+)",
                    "macro_f1": r"(?:macro[_\-\s]*f1|f1[_\-\s]*macro)\s*[:=]\s*([0-9]*\.?[0-9]+)",
                    "balanced_accuracy": r"(?:balanced[_\-\s]*accuracy|balanced[_\-\s]*acc)\s*[:=]\s*([0-9]*\.?[0-9]+)",
                }
                for k, pat in patterns.items():
                    m = re.search(pat, txt, flags=re.I)
                    if m:
                        metrics[k] = float(m.group(1))
                return {"path": str(path), "type": suffix.lstrip("."), "metrics": metrics, "preview": txt[:1000]}
        except Exception as e:
            return {"path": str(path), "error": repr(e)}
        return None

    @staticmethod
    def _flatten_dict(obj: Any, prefix: str = "") -> Dict[str, Any]:
        out = {}
        if isinstance(obj, dict):
            for k, v in obj.items():
                nk = f"{prefix}.{k}" if prefix else str(k)
                out.update(ResultAnalysisAgent._flatten_dict(v, nk))
        else:
            out[prefix] = obj
        return out

    @staticmethod
    def _load_confusion_matrix(path: Path, label_order: List[str]) -> Optional[np.ndarray]:
        try:
            suffix = path.suffix.lower()
            if suffix == ".npy":
                arr = np.load(path)
                return np.asarray(arr, dtype=int)
            if suffix == ".npz":
                data = np.load(path)
                # 优先常见 key
                for key in ["confusion_matrix", "cm", "matrix", "arr_0"]:
                    if key in data:
                        return np.asarray(data[key], dtype=int)
                first = list(data.keys())[0]
                return np.asarray(data[first], dtype=int)
            if suffix == ".csv":
                df = read_csv_robust(path)
                # 情况1：第一列是行标签，列名是类别
                cols_lower = [str(c).strip() for c in df.columns]
                label_set = set(label_order)
                candidate_cols = [c for c in df.columns if str(c).strip() in label_set]
                if len(candidate_cols) >= 2:
                    mat = df[candidate_cols].apply(pd.to_numeric, errors="coerce").values
                    return np.nan_to_num(mat).astype(int)
                # 情况2：纯数值矩阵
                mat = df.apply(pd.to_numeric, errors="coerce").dropna(axis=1, how="all").values
                if mat.ndim == 2 and mat.shape[0] >= 2 and mat.shape[1] >= 2:
                    # 如果第一列是索引导致多一列，尝试取最后 n 列
                    n = len(label_order)
                    if mat.shape[1] == n + 1:
                        mat = mat[:, -n:]
                    return np.nan_to_num(mat).astype(int)
            if suffix == ".txt":
                txt = path.read_text(encoding="utf-8", errors="ignore")
                nums = re.findall(r"-?\d+", txt)
                if len(nums) >= 4:
                    vals = np.array(list(map(int, nums)))
                    n = int(round(math.sqrt(len(vals))))
                    if n * n == len(vals):
                        return vals.reshape(n, n)
        except Exception:
            return None
        return None

    @staticmethod
    def _analyze_cm(cm: np.ndarray, label_order: List[str], matrix_name: str = "confusion_matrix") -> Dict[str, Any]:
        cm = np.asarray(cm)
        n = cm.shape[0]
        labels = label_order[:n]
        if len(labels) < n:
            labels = labels + [f"C{i}" for i in range(len(labels), n)]

        total = cm.sum()
        correct = np.trace(cm)
        accuracy = float(correct / total) if total > 0 else np.nan
        row_sum = cm.sum(axis=1)
        col_sum = cm.sum(axis=0)
        recall = np.divide(np.diag(cm), row_sum, out=np.zeros(n, dtype=float), where=row_sum != 0)
        precision = np.divide(np.diag(cm), col_sum, out=np.zeros(n, dtype=float), where=col_sum != 0)
        f1 = np.divide(2 * precision * recall, precision + recall, out=np.zeros(n, dtype=float), where=(precision + recall) != 0)

        per_class = []
        for i, lab in enumerate(labels):
            per_class.append({
                "class": lab,
                "support": int(row_sum[i]),
                "precision": round(float(precision[i]), 4),
                "recall": round(float(recall[i]), 4),
                "f1": round(float(f1[i]), 4),
                "correct": int(cm[i, i]),
            })

        # 找主要混淆对：真实 i 被预测 j
        confusions = []
        for i in range(n):
            for j in range(n):
                if i == j:
                    continue
                if cm[i, j] > 0:
                    confusions.append({
                        "true": labels[i],
                        "pred": labels[j],
                        "count": int(cm[i, j]),
                        "rate_in_true_class": round(float(cm[i, j] / row_sum[i]), 4) if row_sum[i] > 0 else 0.0,
                    })
        confusions = sorted(confusions, key=lambda x: (x["count"], x["rate_in_true_class"]), reverse=True)

        weakest = sorted(per_class, key=lambda x: x["recall"])[:2]
        strongest = sorted(per_class, key=lambda x: x["recall"], reverse=True)[:2]
        top_conf = confusions[:3]
        text_summary = (
            f"{matrix_name}：总体 Accuracy={accuracy:.4f}，Macro-F1={np.mean(f1):.4f}。"
            f"召回率较高的类别为 {', '.join([x['class'] for x in strongest])}；"
            f"主要薄弱类别为 {', '.join([x['class'] for x in weakest])}。"
        )
        if top_conf:
            pairs = "; ".join([f"真实{x['true']}→预测{x['pred']} {x['count']}次" for x in top_conf])
            text_summary += f" 主要误分方向：{pairs}。"
        if any((x["true"], x["pred"]) in [("HE", "HA"), ("HA", "HE")] for x in top_conf):
            text_summary += " 这说明 HE 与 HA 的边界仍是高光谱叶片检测中的关键难点，需要继续关注无症状病叶与健康叶的细微光谱-空间差异。"

        return {
            "matrix_name": matrix_name,
            "matrix": cm.astype(int).tolist(),
            "labels": labels,
            "accuracy": round(accuracy, 4) if not np.isnan(accuracy) else None,
            "macro_precision": round(float(np.mean(precision)), 4),
            "macro_recall": round(float(np.mean(recall)), 4),
            "macro_f1": round(float(np.mean(f1)), 4),
            "per_class": per_class,
            "top_confusions": top_conf,
            "text_summary": text_summary,
        }

    @staticmethod
    def _summarize_metrics(metric_files: List[Dict[str, Any]]) -> str:
        best_items = []
        for mf in metric_files:
            br = mf.get("best_row")
            if br:
                model = br.get("model") or br.get("Model") or br.get("name") or Path(mf.get("path", "")).stem
                score_by = mf.get("best_by")
                score = br.get(score_by) if score_by else None
                best_items.append(f"{model} 在 {score_by} 上较优，数值为 {score}")
        if best_items:
            return "指标文件自动筛选结果：" + "；".join(best_items[:5]) + "。"
        return "已读取指标文件，但未能识别统一的模型评分列；建议在 metrics_summary.csv 中包含 model、split、accuracy、macro_f1、balanced_accuracy 等列。"


# -----------------------------
# 4. 预处理策略 Agent
# -----------------------------

class PreprocessStrategyAgent(BaseAgent):
    name = "PreprocessStrategyAgent"

    def run(self, ctx: AgentContext) -> None:
        ctx.log(f"{self.name}: 生成预处理策略与风险提示。")
        audit = ctx.artifacts.get("data_audit", {})
        raw_size = audit.get("raw_dir", {}).get("total_size_bytes", 0) or 0
        processed_size = audit.get("processed_dir", {}).get("total_size_bytes", 0) or 0
        label_dist = audit.get("label_distribution", {}) or {}

        suggestions = []
        if raw_size > 0 and processed_size > 0:
            suggestions.append(
                f"原始数据约 {human_size(raw_size)}，预处理后约 {human_size(processed_size)}，"
                f"压缩/整理比例约 {raw_size / processed_size:.2f}×。建议在汇报中强调：数据底座将原始高光谱文件转换为可训练张量，降低服务器训练 IO 压力。"
            )
        else:
            suggestions.append("未同时检测到 raw_dir 和 processed_dir 的大小，建议补充文件夹大小截图，证明预处理后数据约 40GB。")

        if label_dist:
            total = sum(label_dist.values())
            min_label = min(label_dist, key=label_dist.get)
            max_label = max(label_dist, key=label_dist.get)
            suggestions.append(
                f"当前样本总数 {total}，类别分布为 {label_dist}。"
                f"最多类别为 {max_label}，最少类别为 {min_label}。训练时建议保留 balanced accuracy、macro-F1 和类别召回率，避免只看 Accuracy。"
            )

        suggestions.append("预处理报告建议包含：单叶 ROI 提取、质量控制、黑白校正后反射率保存、PCA/尺寸缩放、NPZ 数据底座生成。")
        suggestions.append("若同一 source_image 中含 2/4 片叶，划分训练/测试时应按 source_image 分组，避免同源图像泄漏。")
        suggestions.append("混淆矩阵分析重点关注 HA 与 HE 的互相误分；ND 和 HS 通常较容易识别，可作为模型是否学到病斑/缺素特征的参照。")

        ctx.artifacts["preprocess_strategy"] = suggestions
        write_json(ctx.out_dir / "preprocess_strategy.json", suggestions)


# -----------------------------
# 5. 报告生成 Agent
# -----------------------------

class ReportAgent(BaseAgent):
    name = "ReportAgent"

    def run(self, ctx: AgentContext) -> None:
        ctx.log(f"{self.name}: 开始生成 Markdown/HTML 报告和组会 PPT 提纲。")
        report_md = self._build_markdown_report(ctx)
        md_path = ctx.out_dir / "HLB_Hyperspectral_Agent_Report.md"
        md_path.write_text(report_md, encoding="utf-8")

        html_path = ctx.out_dir / "HLB_Hyperspectral_Agent_Report.html"
        html_path.write_text(self._markdown_to_simple_html(report_md), encoding="utf-8")

        ppt_outline = self._build_ppt_outline(ctx)
        ppt_outline_path = ctx.out_dir / "weekly_meeting_ppt_outline.md"
        ppt_outline_path.write_text(ppt_outline, encoding="utf-8")

        form_text = self._build_form_text(ctx)
        form_text_path = ctx.out_dir / "ai_agent_usage_form_text.txt"
        form_text_path.write_text(form_text, encoding="utf-8")

        ctx.artifacts["report_markdown"] = str(md_path)
        ctx.artifacts["report_html"] = str(html_path)
        ctx.artifacts["ppt_outline"] = str(ppt_outline_path)
        ctx.artifacts["form_text"] = str(form_text_path)

        # 可选：如果安装了 python-pptx，生成一个简易 PPTX 证明稿
        try:
            pptx_path = self._try_build_pptx(ctx)
            if pptx_path:
                ctx.artifacts["auto_pptx"] = str(pptx_path)
        except Exception as e:
            ctx.log(f"{self.name}: python-pptx 简易 PPT 生成失败，不影响主报告：{repr(e)}")

        write_json(ctx.out_dir / "agent_artifacts.json", ctx.artifacts)
        (ctx.out_dir / "agent_run.log").write_text("\n".join(ctx.logs), encoding="utf-8")
        ctx.log(f"{self.name}: 报告生成完成。")

    @staticmethod
    def _build_markdown_report(ctx: AgentContext) -> str:
        audit = ctx.artifacts.get("data_audit", {})
        result = ctx.artifacts.get("result_analysis", {})
        strategy = ctx.artifacts.get("preprocess_strategy", [])
        lines = []
        lines.append(f"# {ctx.project_name}\n")
        lines.append(f"生成时间：{now_str()}\n")
        lines.append("## 1. 项目定位\n")
        lines.append(
            "本项目构建了一个面向柑橘黄龙病叶片高光谱检测的科研辅助 Agent 流水线，"
            "用于连接数据预处理、数据划分、模型训练结果分析、混淆矩阵解读和组会汇报材料生成。"
            "其核心目标不是替代专业判断，而是提高科研工程流程的可复用性、可追踪性和汇报效率。\n"
        )

        lines.append("## 2. 数据审计\n")
        raw = audit.get("raw_dir", {})
        proc = audit.get("processed_dir", {})
        res = audit.get("results_dir", {})
        data_rows = [
            {"项目": "原始数据目录", "路径": raw.get("path"), "文件数": raw.get("n_files"), "大小": raw.get("total_size_human")},
            {"项目": "预处理数据目录", "路径": proc.get("path"), "文件数": proc.get("n_files"), "大小": proc.get("total_size_human")},
            {"项目": "训练结果目录", "路径": res.get("path"), "文件数": res.get("n_files"), "大小": res.get("total_size_human")},
        ]
        lines.append(markdown_table(pd.DataFrame(data_rows)))
        ratio = audit.get("compression_ratio_raw_to_processed")
        reduction = audit.get("space_reduction_percent")
        if ratio:
            lines.append(f"\n预处理后相对于原始数据的整理/压缩比例约为 **{ratio:.2f}×**，空间降低约 **{reduction:.2f}%**。\n")
        else:
            lines.append("\n未能同时计算原始数据与预处理数据大小，建议补充文件夹大小截图。\n")

        label_dist = audit.get("label_distribution", {})
        if label_dist:
            df_label = pd.DataFrame([{"类别": k, "数量": v} for k, v in label_dist.items()])
            lines.append("### 类别分布\n")
            lines.append(markdown_table(df_label))
            lines.append("\n类别分布用于判断是否存在样本不均衡，并决定是否需要重点报告 Macro-F1、Balanced Accuracy 和 per-class recall。\n")

        lines.append("## 3. Agent 模块与核心流程\n")
        module_df = pd.DataFrame([
            {"Agent": "DataAuditAgent", "功能": "扫描原始数据、预处理数据和结果目录，统计文件数量、大小和类别分布"},
            {"Agent": "PreprocessStrategyAgent", "功能": "根据数据规模、类别分布和高光谱任务特点生成预处理策略与风险提示"},
            {"Agent": "SplitDesignAgent", "功能": "生成 5-fold OOF 与 6:2:2 划分，优先按 source_image 分组以降低数据泄漏"},
            {"Agent": "ResultAnalysisAgent", "功能": "解析 metrics、classification report 和混淆矩阵，自动总结模型优劣与主要误分方向"},
            {"Agent": "ReportAgent", "功能": "生成报告、课程表单说明和组会 PPT 提纲"},
        ])
        lines.append(markdown_table(module_df))

        lines.append("\n## 4. 预处理策略总结\n")
        for s in strategy:
            lines.append(f"- {s}")

        lines.append("\n## 5. 模型结果与混淆矩阵分析\n")
        summaries = result.get("analysis_summary", [])
        if summaries:
            for s in summaries:
                lines.append(f"- {s}")
        else:
            lines.append("- 当前未读取到可解析的指标文件或混淆矩阵文件。建议在 results_dir 中保存 metrics_summary.csv、oof_confusion_matrix.csv 和 test_confusion_matrix.csv。")

        cms = result.get("confusion_matrices", [])
        for cmr in cms:
            lines.append(f"\n### {cmr.get('matrix_name')}\n")
            labels = cmr.get("labels", [])
            mat = cmr.get("matrix", [])
            if labels and mat:
                df_cm = pd.DataFrame(mat, index=[f"True_{x}" for x in labels], columns=[f"Pred_{x}" for x in labels])
                lines.append(markdown_table(df_cm.reset_index().rename(columns={"index": "真实/预测"})))
            per = cmr.get("per_class", [])
            if per:
                lines.append("\n类别级指标：\n")
                lines.append(markdown_table(pd.DataFrame(per)))
            conf = cmr.get("top_confusions", [])
            if conf:
                lines.append("\n主要误分方向：\n")
                lines.append(markdown_table(pd.DataFrame(conf)))

        imgs = result.get("image_confusion_matrices", [])
        if imgs:
            lines.append("\n### 已检测到的图片版混淆矩阵\n")
            for img in imgs:
                copied = img.get("copied_to")
                if copied:
                    lines.append(f"- {Path(copied).name}：已复制到 figures 目录。建议同时导出 CSV 版本便于自动计算。")

        lines.append("\n## 6. 后续优化建议\n")
        lines.append("1. 保留 5-fold OOF 作为模型稳定性评估，6:2:2 作为最终独立测试展示。")
        lines.append("2. 汇报中明确说明是否按 source_image 分组，避免同一原始图像中的多片叶进入不同集合。")
        lines.append("3. 除 Accuracy 外，重点展示 Macro-F1、Balanced Accuracy、HA/HE 的召回率和混淆矩阵。")
        lines.append("4. 若 HA 与 HE 仍严重混淆，后续可以尝试注意力模块、特征波段选择、空间-光谱联合输入和更严格的优质样本筛选。")
        lines.append("5. 将每次服务器实验的配置、日志、权重、指标和图像统一保存，形成可追溯实验档案。")

        lines.append("\n## 7. 可作为课程/申报证明材料的文件\n")
        lines.append("- data_audit.json：数据规模、类别分布和目录扫描结果。")
        lines.append("- split_5fold_oof.csv / split_6_2_2.csv：实验划分方案。")
        lines.append("- result_analysis.json：指标和混淆矩阵自动分析。")
        lines.append("- HLB_Hyperspectral_Agent_Report.html：自动生成的 Agent 流水线报告。")
        lines.append("- weekly_meeting_ppt_outline.md：可直接转化为组会 PPT 的提纲。")
        return "\n".join(lines)

    @staticmethod
    def _markdown_to_simple_html(md: str) -> str:
        # 为保证无额外依赖，这里做一个很简单的 HTML 包装，不做完整 Markdown 解析。
        escaped = (
            md.replace("&", "&amp;")
              .replace("<", "&lt;")
              .replace(">", "&gt;")
        )
        # 粗略处理标题和换行
        html_lines = []
        for line in escaped.splitlines():
            if line.startswith("# "):
                html_lines.append(f"<h1>{line[2:]}</h1>")
            elif line.startswith("## "):
                html_lines.append(f"<h2>{line[3:]}</h2>")
            elif line.startswith("### "):
                html_lines.append(f"<h3>{line[4:]}</h3>")
            else:
                html_lines.append(f"<pre>{line}</pre>" if line.startswith("|") else f"<p>{line}</p>")
        return f"""<!doctype html>
<html lang="zh-CN">
<head>
<meta charset="utf-8">
<title>HLB Hyperspectral Agent Report</title>
<style>
body {{ font-family: Arial, 'Microsoft YaHei', sans-serif; margin: 40px; line-height: 1.6; color: #222; }}
h1 {{ color: #0b5a32; }}
h2 {{ color: #176b45; border-bottom: 1px solid #ddd; padding-bottom: 4px; }}
pre {{ background: #fafafa; padding: 4px 8px; white-space: pre-wrap; font-family: Consolas, monospace; }}
p {{ margin: 6px 0; }}
</style>
</head>
<body>
{''.join(html_lines)}
</body></html>"""

    @staticmethod
    def _build_ppt_outline(ctx: AgentContext) -> str:
        audit = ctx.artifacts.get("data_audit", {})
        result = ctx.artifacts.get("result_analysis", {})
        raw_size = audit.get("raw_dir", {}).get("total_size_human", "NA")
        proc_size = audit.get("processed_dir", {}).get("total_size_human", "NA")
        label_dist = audit.get("label_distribution", {})
        summaries = result.get("analysis_summary", [])

        lines = []
        lines.append("# 组会 PPT 提纲：500+片高光谱数据建模实验\n")
        slides = [
            ("封面", "组会汇报 / 基于高光谱图像的柑橘黄龙病叶片检测 / 汇报人 / 日期"),
            ("本周完成的任务", "完成 500 余片优质叶片数据预处理，构建可用于深度学习训练的数据底座；完成 5-fold OOF 和 6:2:2 两套评估实验。"),
            ("预处理与数据压缩", f"原始数据：{raw_size}；预处理后：{proc_size}。重点说明 ROI 提取、质量控制、NPZ 保存、训练 IO 压力降低。"),
            ("样本组成", f"类别分布：{label_dist}。说明 HA/HE/ND/HS 的样本数量和挑选标准。"),
            ("实验设计", "5-fold OOF 用于整体稳定性评估；6:2:2 用于最终独立测试。若使用 source_image 分组，强调避免数据泄漏。"),
            ("5-fold OOF 结果", "放置 oof_confusion_matrix 和 OOF 指标表，分析各类别召回率。"),
            ("6:2:2 测试集结果", "放置 test_confusion_matrix 和测试集指标表，重点分析最终泛化能力。"),
            ("主要问题分析", "围绕 HE 与 HA 的互相混淆展开，说明无症状病叶与健康叶在光谱-空间特征上相近。"),
            ("总结与后续计划", "继续优化模型结构、注意力模块、特征波段选择、样本质量控制和更严格的独立测试。"),
            ("致谢", "感谢老师和师兄们批评与指导。"),
        ]
        for i, (title, body) in enumerate(slides, 1):
            lines.append(f"## Slide {i}: {title}\n{body}\n")
        if summaries:
            lines.append("## 可直接放入 PPT 的自动分析句\n")
            for s in summaries:
                lines.append(f"- {s}")
        return "\n".join(lines)

    @staticmethod
    def _build_form_text(ctx: AgentContext) -> str:
        audit = ctx.artifacts.get("data_audit", {})
        result = ctx.artifacts.get("result_analysis", {})
        raw_size = audit.get("raw_dir", {}).get("total_size_human", "NA")
        proc_size = audit.get("processed_dir", {}).get("total_size_human", "NA")
        label_dist = audit.get("label_distribution", {})
        summaries = result.get("analysis_summary", [])
        summary_text = "；".join(summaries[:3]) if summaries else "系统能够读取训练指标和混淆矩阵，并自动生成结果分析与汇报提纲。"

        return f"""我构建了一个面向柑橘黄龙病叶片检测研究的高光谱科研辅助 Agent。该 Agent 服务于我的电子信息专业研究方向，即基于高光谱图像的柑橘黄龙病叶片检测与分类。项目的核心痛点是：原始高光谱数据体积大、单叶 ROI 提取与质量控制复杂、训练实验版本多、5-fold OOF 与 6:2:2 结果需要快速比较和解释。为解决这些问题，我将 Agent 引入科研流程中，辅助完成数据审计、预处理结果记录、实验划分设计、模型结果分析和组会汇报材料生成。

在数据层面，Agent 会自动扫描原始数据目录、预处理数据目录和训练结果目录，统计文件数量、数据大小和类别分布。本次流程中，原始数据目录大小为 {raw_size}，预处理后数据目录大小为 {proc_size}，类别分布为 {label_dist}。在实验层面，Agent 可以基于 manifest.csv 自动生成 5-fold OOF 和 6:2:2 划分方案，并优先按 source_image 分组，降低同一原始图像中多片叶片进入不同集合造成的数据泄漏风险。在结果层面，Agent 能够读取 metrics_summary、classification report 和混淆矩阵，自动计算 Accuracy、Macro-F1、Precision、Recall 等指标，并总结主要误分方向，例如 HE 与 HA 的混淆、ND 和 HS 的识别稳定性等。当前结果分析摘要为：{summary_text}

该流程已经形成可复用的科研工程流水线，输出包括 data_audit.json、split_5fold_oof.csv、split_6_2_2.csv、result_analysis.json、HTML 报告和组会 PPT 提纲。它提升了高光谱数据管理、模型评估和科研汇报的效率，也使实验过程更加可追溯。"""

    @staticmethod
    def _try_build_pptx(ctx: AgentContext) -> Optional[Path]:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except Exception:
            return None

        pptx_path = ctx.out_dir / "HLB_Hyperspectral_Agent_Proof.pptx"
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        def add_slide(title: str, bullets: Sequence[str]) -> None:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.6))
            tf = title_box.text_frame
            tf.text = title
            tf.paragraphs[0].font.size = Pt(30)
            tf.paragraphs[0].font.bold = True
            body = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.8), Inches(5.6))
            btf = body.text_frame
            btf.clear()
            for i, bullet in enumerate(bullets):
                p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(20)

        audit = ctx.artifacts.get("data_audit", {})
        result = ctx.artifacts.get("result_analysis", {})
        raw_size = audit.get("raw_dir", {}).get("total_size_human", "NA")
        proc_size = audit.get("processed_dir", {}).get("total_size_human", "NA")
        label_dist = audit.get("label_distribution", {})
        summaries = result.get("analysis_summary", [])

        add_slide("AI Agent 辅助构建高光谱黄龙病检测科研流水线", [
            "研究方向：基于高光谱图像的柑橘黄龙病叶片检测",
            "核心任务：数据预处理、实验划分、结果分析与组会汇报自动化",
            "Agent 输出：数据审计、划分文件、结果报告、PPT 提纲",
        ])
        add_slide("数据审计与预处理", [
            f"原始数据目录大小：{raw_size}",
            f"预处理数据目录大小：{proc_size}",
            "预处理目标：单叶 ROI、质量控制、反射率/张量保存、降低训练 IO 压力",
        ])
        add_slide("样本组成与实验划分", [
            f"类别分布：{label_dist}",
            "5-fold OOF：评估模型稳定性",
            "6:2:2：训练/验证/测试独立评估",
            "优先按 source_image 分组，降低数据泄漏风险",
        ])
        add_slide("模型结果自动分析", summaries[:4] if summaries else [
            "读取 metrics_summary 与 confusion_matrix",
            "自动计算 Accuracy、Macro-F1、Precision、Recall",
            "识别主要误分方向，例如 HE 与 HA 的互相混淆",
        ])
        add_slide("总结与后续计划", [
            "继续优化空间-光谱联合模型与注意力模块",
            "重点提升 HA 与 HE 的区分能力",
            "保留完整实验日志，提高科研可复现性",
        ])
        prs.save(pptx_path)
        return pptx_path


# -----------------------------
# Orchestrator
# -----------------------------

class ResearchAgentOrchestrator:
    def __init__(self, agents: Sequence[BaseAgent]):
        self.agents = list(agents)

    def run(self, ctx: AgentContext) -> None:
        ensure_dir(ctx.out_dir)
        ctx.log("ResearchAgentOrchestrator: 启动高光谱黄龙病科研 Agent 流水线。")
        for agent in self.agents:
            try:
                agent.run(ctx)
            except Exception as e:
                ctx.log(f"{agent.name}: 运行异常：{repr(e)}")
        ctx.log("ResearchAgentOrchestrator: 全部流程结束。")


# -----------------------------
# CLI
# -----------------------------

def parse_args(argv: Optional[Sequence[str]] = None) -> argparse.Namespace:
    p = argparse.ArgumentParser(description="HLB Hyperspectral Research Agent Pipeline")
    p.add_argument("--project_name", default="Citrus-HLB Hyperspectral Research Agent", help="项目名称")
    p.add_argument("--manifest", default=None, help="样本清单 CSV，建议包含 sample_id,label,source_image,path")
    p.add_argument("--raw_dir", default=None, help="原始高光谱数据目录")
    p.add_argument("--processed_dir", default=None, help="预处理后 NPZ/ROI 数据目录")
    p.add_argument("--results_dir", default=None, help="训练结果目录，包含 metrics、confusion_matrix 等")
    p.add_argument("--out_dir", required=True, help="Agent 输出目录")
    p.add_argument("--label_col", default="label", help="manifest 中类别列名")
    p.add_argument("--group_col", default="source_image", help="manifest 中分组列名，如 source_image；没有则传空字符串")
    p.add_argument("--id_col", default="sample_id", help="manifest 中样本ID列名")
    p.add_argument("--label_order", default=",".join(DEFAULT_LABEL_ORDER), help="类别顺序，逗号分隔，如 HA,HE,ND,HS")
    p.add_argument("--make_splits", action="store_true", help="是否生成 5-fold OOF 与 6:2:2 划分文件")
    return p.parse_args(argv)


def path_or_none(x: Optional[str]) -> Optional[Path]:
    if x is None:
        return None
    x = str(x).strip()
    if x == "":
        return None
    return Path(x)


def main(argv: Optional[Sequence[str]] = None) -> int:
    args = parse_args(argv)
    label_order = [x.strip() for x in args.label_order.split(",") if x.strip()]
    ctx = AgentContext(
        project_name=args.project_name,
        manifest_path=path_or_none(args.manifest),
        raw_dir=path_or_none(args.raw_dir),
        processed_dir=path_or_none(args.processed_dir),
        results_dir=path_or_none(args.results_dir),
        out_dir=Path(args.out_dir),
        label_col=args.label_col,
        group_col=args.group_col.strip() if args.group_col and args.group_col.strip() else None,
        id_col=args.id_col.strip() if args.id_col and args.id_col.strip() else None,
        label_order=label_order,
        make_splits=bool(args.make_splits),
    )

    agents: List[BaseAgent] = [
        DataAuditAgent(),
        SplitDesignAgent(),
        PreprocessStrategyAgent(),
        ResultAnalysisAgent(),
        ReportAgent(),
    ]
    orchestrator = ResearchAgentOrchestrator(agents)
    orchestrator.run(ctx)

    print("\n=== 输出文件 ===")
    for k, v in ctx.artifacts.items():
        print(f"{k}: {v}")
    print(f"\n完成。请查看输出目录：{ctx.out_dir}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
